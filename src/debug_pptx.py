from pathlib import Path
from pptx import Presentation


PPTX_PATH = Path("data/raw/L1.pptx")


def debug_pptx(path: Path) -> None:
    print(f"File: {path}")
    print(f"Exists: {path.exists()}")

    prs = Presentation(path)
    print(f"Slides: {len(prs.slides)}")

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print()
        print(f"--- SLIDE {slide_idx} ---")
        print(f"Number of shapes: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            shape_type = getattr(shape, "shape_type", None)
            has_text_frame = getattr(shape, "has_text_frame", False)
            text = getattr(shape, "text", "")

            print(f"Shape {shape_idx}: type={shape_type}, has_text_frame={has_text_frame}")

            if text and text.strip():
                print("TEXT:")
                print(text.strip()[:500])


if __name__ == "__main__":
    debug_pptx(PPTX_PATH)