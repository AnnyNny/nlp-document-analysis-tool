from pathlib import Path
import re

import fitz  # PyMuPDF
from pptx import Presentation


RAW_DIR = Path("data/raw")
EXTRACTED_DIR = Path("data/extracted")


def normalize_text(text: str) -> str:
    text = text.replace("\x0b", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text_from_pptx(pptx_path: Path) -> str:
    presentation = Presentation(pptx_path)
    document_parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"--- SLIDE {slide_index} ---"]

        for shape in slide.shapes:
            text = getattr(shape, "text", "")

            if text and text.strip():
                slide_parts.append(normalize_text(text))

            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    cells = []
                    for cell in row.cells:
                        cell_text = normalize_text(cell.text)
                        if cell_text:
                            cells.append(cell_text)
                    if cells:
                        slide_parts.append(" | ".join(cells))

            if hasattr(shape, "shapes"):
                for nested_shape in shape.shapes:
                    nested_text = getattr(nested_shape, "text", "")
                    if nested_text and nested_text.strip():
                        slide_parts.append(normalize_text(nested_text))

        document_parts.append("\n".join(slide_parts))

    return "\n\n".join(document_parts)


def extract_text_from_pdf(pdf_path: Path) -> str:
    document_parts = []

    with fitz.open(pdf_path) as pdf_document:
        for page_index, page in enumerate(pdf_document, start=1):
            page_text = page.get_text("text")
            page_text = normalize_text(page_text)
            document_parts.append(f"--- PAGE {page_index} ---\n{page_text}")

    return "\n\n".join(document_parts)


def extract_text_from_txt(txt_path: Path) -> str:
    return txt_path.read_text(encoding="utf-8", errors="ignore")


def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        return extract_text_from_pptx(file_path)

    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)

    if suffix == ".txt":
        return extract_text_from_txt(file_path)

    raise ValueError(f"Unsupported file type: {file_path.suffix}")


def extract_all_documents(raw_dir: Path = RAW_DIR, output_dir: Path = EXTRACTED_DIR) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)

    supported_files = sorted(
        list(raw_dir.glob("*.pptx"))
        + list(raw_dir.glob("*.pdf"))
        + list(raw_dir.glob("*.txt"))
    )

    if not supported_files:
        print(f"No supported files found in {raw_dir.resolve()}")
        return

    for file_path in supported_files:
        try:
            extracted_text = extract_text_from_file(file_path)
            output_path = output_dir / f"{file_path.stem}.txt"
            output_path.write_text(extracted_text, encoding="utf-8")

            line_count = len(extracted_text.splitlines())
            word_count = len(extracted_text.split())

            print(f"Extracted: {file_path.name} -> {output_path}")
            print(f"  Lines: {line_count}")
            print(f"  Words: {word_count}")

        except Exception as error:
            print(f"Failed to extract {file_path.name}: {error}")


if __name__ == "__main__":
    extract_all_documents()